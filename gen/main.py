import openai




# user input

def get_input():
  
  q5_names_ = []
  q10_info = []
  q4_demographics = {}
  time = ["_", "400", "800", "1200"]
  
  q1 = input("Как называется ваша компания? ")
  q2 = input("Какие услуги вы предоставляете? ")
  q3 = input("На какой рынок ориентируется ваша деятельность? ")
  q4_1 = input("На людей какого ВОЗРАСТА ориентирована деятельность вашего бизнес проекта? ")
  q4_2 = input("На людей какого ПОЛА ориентирована деятельность вашего бизнес проекта? ")
  q4_3 = input("На людей какого РЕГИОНА(местности) проживания ориентирована деятельность вашего бизнес проекта? ")
  q4_4 =input("На людей каких ИНТЕРЕСОВ ориентирована деятельность вашего бизнес проекта? ") 
  q4_demographics = {"age":f'{q4_1}', "gender":f'{q4_2}', "location":f'{q4_3}', "interests":f'{q4_4}'}
  def add_name():
    q5_name = input("Назовите имя и фамилию одного из ваших сотрудников: ")
    q5_xper = input("Опишите его роль в команде: ")
    q5_names_.append({"name":q5_name, "expertise":q5_xper})
  add_name()
  more_names = input("Хотите добавить данные следующего сотрудника?(y/n) ").upper()
  while more_names == "Y":
    add_name()
    more_names = input("Хотите добавить данные следующего сотрудника?(y/n) ").upper()
  q6 = input("Чего ваш проект уже достиг на данном этапе? ")
  q7 = input("На что вы направите будующие инвестиции? ")
  q8 = input("На какой размер инвестиций вы рассчитываете? ")
  q9 = input("К какому сроку вы планируете оправдать инвестиции? ")
  def add_info():
    q10_name = input("Укажите название интернет ресурса/имя и фамилию ответчика ")
    q10_link = input("Введите ссылку или номер телефона ")
    q10_info.append({"name": f'{q10_name}', "link":f'{q10_link}'})
  add_contact = input("Хотите добавить контактную информацию?(y/n) ").upper()
  while add_contact == "Y":
    add_info()
    add_contact = input("Хотите добавить контактную информацию?(y/n) ").upper()
  q11_num = input("Сколько времени должен занять ваш питчдек? Введите циыру от одного до трех (1 - 2 минуты; 2 - 5 минут; 3 - 12 минут) ")
  try:
    q11 = time[q11_num]
  except:
    q11 = time[2]

  answers = (q1, q2, q3, q4_demographics, q5_names_, q6, q7, q8, q9, q10_info, q11)
  return answers


sample_input = {'company_name' : 'SimpleCoffe',
                'business_type' : 'A coffeeshop', 
                'target market': 'Saint-Petersburg, Russia',
                'target audience': {"age":"all ages", "gender":"all genders", "location":"Saint-Petersburg", "interests":"drinking tasty coffe"},
                'project team': [{"name":"Daria Ivanova",  "expertise":"Barista"}, {"name":"Dmitry Petrov", "expertise": "accountant"}],
                'The project have already achieved': "monthly proffit of 25000 rubles, 30000 visitors per month",
                'investment target': 'aquiring new equipment, making better product, doubbling our revenue',
                'investment amount': '10 million rubles', 
                'deadline on investment promises': '1 year',
                'contact information': "connect to us directly using the number 123456789",
                'length(symbols)': '200'} 

sample_input_rus = {'company_name' : "ПростоКофе",
                'business_type' : "кофейня", 
                'target market': "Санкт-Петербург, россия",
                'target audience': {"age":"всех возростов", "gender":"любого пола", "location":"Санкт-Петербург", "interests":"Разнообразных интересов"},
                'project team': [{"name":"Дарья Иванова",  "expertise":"Бариста"}, {"name":"Дмитрий Петров", "expertise": "Бухгалтер"}],
                'The project have already achieved': "20000 посетителей в месяц, прибыль более полумиллиона рублей",
                'investment target': 'aquiring new equipment, making better product, doubbling our revenue',
                'investment amount': '10 миллионов рублей', 
                'deadline on investment promises': '1 год',
                'contact information': "", 
                'length(symbols)': '200'}   

def formated_input():
  answers = get_input()
  user_input = {'company name' : f'{answers[0]}',
                  'business type' : f'{answers[1]}', 
                  'target market': f'{answers[2]}',
                  'target audience': answers[3],
                  'project team': answers[4],
                  'The project have already achieved': f"{answers[5]}",
                  'investment target': f'{answers[6]}',
                  'investment amount': f'{answers[7]}', 
                  'deadline on investment promises': f'{answers[8]}',
                  'information': f'{answers[9]}',
                  'length(symbols)': f'{answers[10]}'} 
  return user_input


def compress_dict(dict_):
  compressed = ''
  for key, data in dict_.items():
    compressed += f' {key}: {data},'
  return compressed

def get_txt_input(dict_):
  final_str = ''
  for key, data in dict_.items():
    if type(data) == list:
      temp_str = ''
      for subdict_ in data:
        temp_str += compress_dict(subdict_)
      final_str += f'{key}: {temp_str}\n'
    elif type(data) == dict:
      final_str += f'{key}: {compress_dict(data)}\n'
    else:
      final_str += f'{key}: {data}\n'
  return final_str

#print(get_txt_input(sample_input_rus))
#print(get_txt_input(sample_input))

# requests



#completion = openai.ChatCompletion.create(
# model="gpt-3.5-turbo", 
# messages=[{"role": "user", "content": "What is the OpenAI mission?"}]
#)

#translation = openai.ChatCompletion.create(
#  model = "gpt-3.5-turbo", 
#  messages=[{"role": "user", "content": f"translate to english: {get_txt_input(sample_input_rus)}"}]
#)

import json

def get_json():
  
  with open("gen/key.txt") as f:
    openai.api_key = f.read()

  with open("gen/questions.txt") as f:
    question = f.read()

  test_resp = openai.ChatCompletion.create(
    model = "gpt-3.5-turbo", 
    messages=[
      {"role": "system", "content": "You are a helpful assistant."},
      {"role": "user", "content": f"{question}\n {get_txt_input(formated_input())}"}
      ]
  )

  contents = test_resp["choices"][0]["message"]["content"]

  print(contents)

  with open("gen/response_log.json", "+w") as f:
    f.write(contents)

  response = ''
  
  with open("gen/response_log.json") as f:
    response = json.loads(f.read())

  return response

# translation

#import goslate

#gs = goslate.Goslate
#def translate(path):
#  with open(path) as f:
#    response = json.dumps(f.read())
#    response = gs.translate(response, 'en', 'rus' )
# return response

#print(translate("gen/response_log.json"))
# get slides

import collections
import collections.abc
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

a, b = ("THIS IS A TEST", "this is a test")

prs = Presentation()
slide_layout = prs.slide_layouts[6]

def get_title_slides():
  global prs, slide_layout
  response = get_json()
  for element in response["slides"]:
    title_text = element["title"]
    discr = get_txt_input(element["content"])
    count = 0
    temp = ''
    string = ''
    for symbol in discr:
      if count == 0 and symbol == ' ':
        symbol = ''   
      if count >= 50 or symbol == ':':
        temp += symbol
        if (symbol in [' ', ':']):
          string += f'{temp}\n'
          temp = ''
          count = 0
        else:
          count += 1
      else:
        temp += symbol
        count += 1
    string += temp
    
    discr = string.replace('_', ' ').replace('.,', '.')

    size1 = 60
    if len(title_text) > 20:
      size1 = 41


    slide = prs.slides.add_slide(slide_layout)
    
    left = Inches(5)
    top = Inches(1)
    width = Inches(0.1)
    height = Inches(0.1)
    txBox = slide.shapes.add_textbox(left, top, 0, 0)
    txBox.word_wrap = True
    title_box = txBox.text_frame
    title_par = title_box.add_paragraph()
    title_par.word_wrap = True
    title_par.alignment = PP_ALIGN.LEFT
    title_par.text = title_text
 
    title_par.font.bold = True
    title_par.font.size = Pt(size1) 
    
    discr_par = title_box.add_paragraph()
    discr_par.alignment = PP_ALIGN.LEFT
    discr_par.text = discr
    discr_par.font.size = Pt(20) 

  

 
  prs.save('gen/test.pptx')

get_title_slides()



